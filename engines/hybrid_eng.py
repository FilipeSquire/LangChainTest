from dotenv import load_dotenv
from unstructured.partition.pdf import partition_pdf
from langchain_groq import ChatGroq
from langchain_core.prompts import ChatPromptTemplate
from langchain_openai import ChatOpenAI
from langchain_core.output_parsers import StrOutputParser
import uuid
from langchain.vectorstores import Chroma
from langchain.storage import InMemoryStore
from langchain_core.documents import Document
from langchain_openai import OpenAIEmbeddings
from langchain.retrievers.multi_vector import MultiVectorRetriever
from langchain.retrievers import BM25Retriever, EnsembleRetriever
from langchain_core.runnables import RunnablePassthrough, RunnableLambda
from langchain_core.messages import SystemMessage, HumanMessage
from base64 import b64decode
from engines.prompts import system_finance_prompt
from base64 import b64decode, b64decode as _b64
import re
import time
import io
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from pptx.util import Pt
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typing import Any, Dict, List, Optional

load_dotenv()

class HybridEngine():

    def __init__(self, pdf):
        #Preprocessing variables
        self.pdf = pdf
        self.chunks = None
        self.tables, self.texts, self.images = [], [], []
        self.text_summaries, self.table_summaries = None, None

        #Vector store variables
        self.vectorstore = Chroma(collection_name='multi_modal_rag', embedding_function=OpenAIEmbeddings())
        self.store = InMemoryStore()
        self.id_key = 'doc_id'
        self.dense_retriever = MultiVectorRetriever(
            vectorstore=self.vectorstore,
            docstore=self.store,
            id_key=self.id_key,
            search_kwargs={"k": 30},          # top-40 neighbours
        )

        # Model
        self._hydrid = None

        # Caller
        self.chain = None
        self.chain_with_sources = None



    def _unstructured(self):
        
        # UNSTRUCTED BLOCK
        chunks = partition_pdf(
            file= self.pdf,
            infer_table_structure=True, #extracting table
            strategy = 'hi_res', #mandatory to infer table

            extract_image_block_types=['Image'], #add 'Table' to list to extract image of tables
            # image_output_dir_path = output_path, #if None, images and tables will be saved as base64

            extract_image_block_to_payload=True, #if true, extract base64 for API usage

            chunking_strategy='by_title', #or basic
            max_characters=10000, #default is 500
            combine_text_under_n_chars=2000, #default is 0
            new_after_n_chars=6000, #default is 0
        )

        self.chunks = chunks

        # ELEMENTS SPLIT
        tables, texts, images= [], [], []

        for chunk in chunks:
            if "CompositeElement" in str(type(chunk)):
                chunk_els = chunk.metadata.orig_elements
                for el in chunk_els:
                    if "Table" in str(type(el)):
                        tables.append(el)
                    elif "Image" in str(type(el)):
                        images.append(el)
                    else:
                        texts.append(el)

        self.tables, self.texts, self.images = tables, texts, images
        print('Finished unstructured')


    def _summarization(self):
        # SUMMARIZATION
        prompt_text = """
        You are an especialist in corporate finance tasked with summarizing the text, tables and images.
        Give a concise summary of the table, text or image.
        The tables will be received in format html. Transform this format in order to interpret the table.

        The summary must take special attention to financial-related numbers and statistics, such as monthly or yearly comparisons, debt/loan information, and other subjects related.
        The summary must contain the numerical information of debt, loan, revenue, deficit, and other related topics.
        Always mention in the summary from which of the blocks the content being summarized is part of (Introduction Table, Business Overview, Revenue Split, Key Stakeholders Table, Financial Highlights, Capital Structure)
        Response only with the summary, no additional comment. 
        Do not start your message by saying "Here is a summary" or anything like that. 
        Just give the summart as it is.

        Table or text chunk of Tour Partner Groups: {element}
        """

        prompt = ChatPromptTemplate.from_template(prompt_text)

        model = ChatOpenAI(model="gpt-4o-mini", temperature=0.2,)
        # model = ChatGroq(temperature=0.5, model='llama-3.1-8b-instant')
        summarize_chain = {"element": lambda x: x} | prompt | model | StrOutputParser()

        self.text_summaries = summarize_chain.batch(self.texts, {'max_concurrency':3})
        self.table_summaries = summarize_chain.batch(self.tables, {'max_concurrency':3})
        print('Finished Summarization')


    def _store_load(self):
        #Loading values

        doc_ids = [str(uuid.uuid4()) for _ in self.texts]
        summary_texts = [
            Document(page_content=summary, metadata={self.id_key: doc_ids[i]}) for i, summary in enumerate(self.text_summaries)
        ]
        self.dense_retriever.vectorstore.add_documents(summary_texts)
        self.dense_retriever.docstore.mset(list(zip(doc_ids, self.texts)))

        # Add tables
        tables_ids = [str(uuid.uuid4()) for _ in self.tables]
        summary_tables = [
            Document(page_content=summary, metadata={self.id_key: tables_ids[i]}) for i, summary in enumerate(self.table_summaries)
        ]
        self.dense_retriever.vectorstore.add_documents(summary_tables)
        self.dense_retriever.docstore.mset(list(zip(tables_ids, self.tables)))
        print('Finished store load')

    
    def _hydra(self):

        # lexical (sparse) – BM25 over the same documents
        keys = list(self.store.yield_keys())      # ['id1', 'id2', ...]
        raw_items = self.store.mget(list(self.store.yield_keys()))

        all_docs = [
            item if isinstance(item, Document)
            else Document(page_content=getattr(item, "text", str(item)))
            for item in raw_items
        ]

        bm25 = BM25Retriever.from_documents(all_docs)
        bm25.k = 20                                            # top-20 keyword hits

        # ------------------------------------------------------------------ #
        # 2.  combine them                                                   #
        # ------------------------------------------------------------------ #

        self.hybrid = EnsembleRetriever(
            retrievers=[bm25, self.dense_retriever],
            weights=[0.4, 0.6],               # adjust to bias dense vs. sparse
        )
        print('Finished hydra')


# ---------------------------- RAG PIPE ----------------------------
    def _looks_like_b64(self, s: str) -> bool:
        _B64_RE = re.compile(r"^[A-Za-z0-9+/]+={0,2}$")

        if len(s) < 40 or not _B64_RE.fullmatch(s.strip()):
            return False
        try:
            _b64(s, validate=True)
            return True
        except Exception:
            return False
        
    def _to_str(self, obj) -> str:
        if isinstance(obj, str):
            return obj
        if hasattr(obj, "page_content"):
            return obj.page_content  # LangChain Document
        if hasattr(obj, "text"):
            return obj.text          # unstructured element
        return str(obj)
    
    def _parse_docs(self, docs):
        images, texts = [], []
        for d in docs:
            payload = self._to_str(d)
            if self._looks_like_b64(payload):
                images.append(payload)
            else:
                texts.append(payload)
        return {"images": images, "texts": texts}

    def _build_prompt_two(self, kwargs) -> ChatPromptTemplate:
        ctx      = kwargs["context"]
        question = kwargs["question"]

        parts = [
            {
                "type": "text",
                "text": f"Context:{ctx.get('texts', [])}, Question: {question}",
            }
        ]

        for b64 in ctx.get("images", []):
            parts.append({
                "type": "image_url",
                "image_url": {"url": f"data:image/jpeg;base64,{b64}"},
            })

        messages = [
            SystemMessage(content=system_finance_prompt),
            HumanMessage(content=parts),
        ]
        return ChatPromptTemplate.from_messages(messages)


    def _RAG(self):

        self.chain = (
            {
                'context': self.hybrid | RunnableLambda(self._parse_docs),
                'question': RunnablePassthrough(),
            }
            | RunnableLambda(self._build_prompt_two)
            | ChatOpenAI(model='o3')
            | StrOutputParser()
        )

        self.chain_with_sources = {
            'context': self.hybrid | RunnableLambda(self._parse_docs),
            'question': RunnablePassthrough(),
        } | RunnablePassthrough().assign(
            response=(
                RunnableLambda(self._build_prompt_two)
                | ChatOpenAI(model='o3')
                | StrOutputParser()
            )
        )

        # BETTER FACTOR THE STORE
        # 1. pull every key currently in the store
        keys = list(self.store.yield_keys())                 # ['uuid1', 'uuid2', …]
        # 2. map any raw element → Document(page_content=str)
        fixed_pairs = []
        for k in keys:
            val = self.store.mget([k])[0]                    # fetch single object
            if isinstance(val, Document):
                fixed_pairs.append((k, val))            # already okay
            else:                                       # Title / NarrativeText / Table
                fixed_pairs.append(
                    (k, Document(page_content=getattr(val, "text", str(val))))
                )

        # 3. overwrite the store with the cleaned documents
        self.store.mset(fixed_pairs)
        print('Finished RAG')


    def _generate_pdf(self, text):
        buf = io.BytesIO()
        doc = SimpleDocTemplate(buf, pagesize=letter)

        styles = getSampleStyleSheet()
        body = styles["BodyText"]

        story = []
        for line in text.split("\n"):
            story.append(Paragraph(line or " ", body))
            story.append(Spacer(1, 4))

        doc.build(story)
        buf.seek(0)
        

        return buf.read()
    
    def _replace_placeholder_text(self, file_path: str,placeholder: str,replacement: str,output_path: None) -> None:
        """
        Replace occurrences of `placeholder` in text elements with `replacement`.

        Modifies the presentation and saves to `output_path` or overwrites original.
        """
        try:
            prs = Presentation(file_path)
        except Exception as e:
            raise IOError(f"Unable to open file {file_path}: {e}")

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if placeholder in run.text:
                            run.text = run.text.replace(placeholder, replacement)
                            run.font.size = Pt(5)

        save_path = output_path or file_path
        try:
            prs.save(save_path)
        except Exception as e:
            raise IOError(f"Unable to save updated file to {save_path}: {e}")
    
    def _generate_ppt(self, response):

        ppt_path = './LangChainTest/output/try2.pptx'
        pdf_context = response

        my_prompt = f"""
        The following text contains a series of text blocks that are separated by multiples - like:
        --------------------------------------------------------------------------------------------------------------------

        Bring me back only the code block relative to Block 1., and give me all the information in it in a single line but separated by |

        The text is:

        {pdf_context}
        """
        ppt_response = self.chain_with_sources.invoke(my_prompt)
        self._replace_placeholder_text('./LangChainTest/backupppt.pptx', 'Company Snapshot', ppt_response['response'], ppt_path)


        my_prompt = f"""
        The following text contains a series of text blocks that are separated by multiples - like:
        --------------------------------------------------------------------------------------------------------------------

        Bring me back only the code block relative to Block 2., and give me back exactly as is in the text.

        The text is:

        {pdf_context}
        """
        ppt_response = self.chain_with_sources.invoke(my_prompt)
        self._replace_placeholder_text(ppt_path, 'Business Overview Text', ppt_response['response'],ppt_path)


        my_prompt = f"""
        The following text contains a series of text blocks that are separated by multiples - like:
        --------------------------------------------------------------------------------------------------------------------

        Bring me back only the code block relative to Block 3., and give me back exactly as is in the text.

        The text is:

        {pdf_context}
        """
        ppt_response = self.chain_with_sources.invoke(my_prompt)
        self._replace_placeholder_text(ppt_path, 'Revenue Split Text', ppt_response['response'],ppt_path)


        my_prompt = f"""
        The following text contains a series of text blocks that are separated by multiples - like:
        --------------------------------------------------------------------------------------------------------------------

        Bring me back only the code block relative to Block 4., and give me back exactly as is in the text.

        The text is:

        {pdf_context}
        """
        ppt_response = self.chain_with_sources.invoke(my_prompt)
        self._replace_placeholder_text(ppt_path, 'Key Stakeholders Table', ppt_response['response'],ppt_path)


        my_prompt = f"""
        The following text contains a series of text blocks that are separated by multiples - like:
        --------------------------------------------------------------------------------------------------------------------

        Find the block 5., analyze all metric data and create a table where it is possible to clearly identify the metric and its values.

        Return nothing but the table.

        The text is:

        {pdf_context}
        """
        ppt_response = self.chain_with_sources.invoke(my_prompt)
        self._replace_placeholder_text(ppt_path, 'Financial Highlights Table', ppt_response['response'],ppt_path)


        my_prompt = f"""
        The following text contains a series of text blocks that are separated by multiples - like:
        --------------------------------------------------------------------------------------------------------------------

        Find the block 5., and find within it something related to Commentary. 
        Analyze all the points that are separated in a numeric list, and create a summary. 

        Return nothing but the summary created.

        The text is:

        {pdf_context}
        """
        ppt_response = self.chain_with_sources.invoke(my_prompt)
        self._replace_placeholder_text(ppt_path, 'Financial Highlights Commentary', ppt_response['response'],ppt_path)


        my_prompt = f"""
        The following text contains a series of text blocks that are separated by multiples - like:
        --------------------------------------------------------------------------------------------------------------------

        Find the block 6., analyze all metric data and create a table where it is possible to clearly identify the metric and its values.

        Return nothing but the table.

        The text is:

        {pdf_context}
        """
        ppt_response = self.chain_with_sources.invoke(my_prompt)
        self._replace_placeholder_text(ppt_path, 'Capital Structure Table', ppt_response['response'],ppt_path)


        my_prompt = f"""
        The following text contains a series of text blocks that are separated by multiples - like:
        --------------------------------------------------------------------------------------------------------------------

        Find the block 6., and find within it something related to Capital-structure commentary. 
        Analyze all the points that are separated in a numeric list, and create a summary. 

        Return nothing but the summary created.

        The text is:

        {pdf_context}
        """
        ppt_response = self.chain_with_sources.invoke(my_prompt)
        self._replace_placeholder_text(ppt_path, 'Capital Structure Commentary', ppt_response['response'],ppt_path)

            
    def create_profile(self):

        my_prompt2 = 'Make the company profile. Use the context given'
        i=0
        while True:
            i+=1
            tic = time.perf_counter()
            response = self.chain_with_sources.invoke(my_prompt2)
            print(f'Try {i}')

            if time.perf_counter() - tic >= 60:
                break
        
        # Creating pdf file
        pdf = self._generate_pdf(response['response'])
        self._generate_ppt(response['response'])

        return pdf
    
    def main(self):
        self._unstructured()
        self._summarization()
        self._store_load()
        self._hydra()
        self._RAG()
        print('Finished pipeline')